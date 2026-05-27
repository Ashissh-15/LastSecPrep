from pptx import Presentation
import os

from loaders.image_extractor import (
    extract_image_description
)


def load_pptx(file_path):

    documents=[]

    prs=Presentation(
        file_path
    )


    for slide_number,slide in enumerate(
        prs.slides
    ):

        slide_text=""


        for shape in slide.shapes:

            try:

                # ---------------------
                # Normal text
                # ---------------------

                if hasattr(
                    shape,
                    "text"
                ):

                    text=shape.text.strip()

                    if text:

                        slide_text += (
                            text + "\n"
                        )


                # ---------------------
                # Tables
                # ---------------------

                if shape.has_table:

                    table=shape.table

                    header=None

                    for row_num,row in enumerate(
                        table.rows
                    ):

                        row_data=[]

                        for cell in row.cells:

                            row_data.append(
                                cell.text.strip()
                            )


                        if row_num==0:

                            header=row_data

                            continue


                        if header:

                            row_text=[]

                            for i in range(
                                min(
                                    len(header),
                                    len(row_data)
                                )
                            ):

                                row_text.append(

f"{header[i]} : {row_data[i]}"

                                )

                            slide_text += (
                                " | ".join(
                                    row_text
                                )
                                + "\n"
                            )


                # ---------------------
                # Images
                # ---------------------

                if (
                    shape.shape_type==13
                ):

                    image=shape.image

                    temp_name=(
f"ppt_img_{slide_number}.png"
                    )

                    with open(
                        temp_name,
                        "wb"
                    ) as f:

                        f.write(
                            image.blob
                        )


                    description=(
                        extract_image_description(
                            temp_name
                        )
                    )


                    slide_text += (
                        "\n"
                        + description
                        + "\n"
                    )


                    os.remove(
                        temp_name
                    )


            except Exception as e:

                print(
f"Error reading shape: {e}"
                )


        if slide_text.strip():

            documents.append({

                "text":slide_text,

                "source":file_path,

                "page":
                slide_number+1
            })


    return documents